import streamlit as st
import requests
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches

# Function to fetch multiple images from Unsplash
def fetch_unsplash_images(query, count=5):
    UNSPLASH_ACCESS_KEY = "ybGsMYkr19IKVmanvnpZbihjVUNa8wj4fCXO6VJeuIc"  # Replace with your Unsplash Access Key
    url = f"https://api.unsplash.com/search/photos?query={query}&client_id={UNSPLASH_ACCESS_KEY}&per_page={count}"
    response = requests.get(url)

    if response.status_code == 200:
        data = response.json()
        if data['results']:
            return [result['urls']['regular'] for result in data['results']]
        else:
            return None
    else:
        st.error(f"Error fetching images: {response.status_code}")
        return None

# Function to generate slide content using Google's Generative AI
def generate_slide_content(topic, gemini_api_key, num_slides):
    genai.configure(api_key=gemini_api_key)
    role_prompt = f"""
    You are a presentation maker AI. Your task is to create content details like definition, types, use, benefits, and examples of a PowerPoint presentation.
    The content should be based on the given topic.

    For each slide, provide:
    - A clear slide title (without the word 'Title').
    - 3-4 Key points formatted without bullet characters (*, -, etc.).

    The number of slides should be {num_slides}.

    Topic: {topic}
    """
    model = genai.GenerativeModel("gemini-1.5-flash")
    response = model.generate_content(role_prompt)
    return response.text

# Function to create a PowerPoint presentation
def create_ppt(topic, ai_generated_content, image_urls):
    presentation = Presentation()

    # Add a title slide
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = topic
    subtitle.text = "Generated by Umar_PPT_AI"

    # Parse AI-generated content and add slides
    slides = ai_generated_content.strip().split("Slide")[1:]  # Split content into slides
    for i, slide_content in enumerate(slides[:len(image_urls)]):  # Limit the number of slides based on user input
        lines = slide_content.strip().split("\n")
        slide_title = lines[0].replace("Title:", "").strip()  # Remove "Title:" if present
        bullet_points = [line.lstrip("*- ") for line in lines[1:] if line.strip()]  # Remove * and - from bullet points

        # Add a content slide
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        title.text = slide_title
        content.text = "\n".join(bullet_points)

        # Add an image to the right side of the slide if available
        if i < len(image_urls):
            try:
                image_response = requests.get(image_urls[i])
                if image_response.status_code == 200:
                    with open(f"temp_image_{i}.jpg", "wb") as f:
                        f.write(image_response.content)

                    # Set image position (Right side)
                    left = Inches(5.5)  # Move image to the right
                    top = Inches(1.5)   # Align with content
                    height = Inches(2)  # Reduce size for better layout

                    slide.shapes.add_picture(f"temp_image_{i}.jpg", left, top, height=height)
            except Exception as e:
                st.error(f"Error adding image to slide: {e}")

    # Save the presentation
    ppt_filename = f"{topic.replace(' ', '_')}.pptx"
    presentation.save(ppt_filename)
    return ppt_filename

# Streamlit app
def app():
    st.markdown("""
        <style>
            .stApp {
                background-color:rgb(190, 220, 251);  /* Deep Seek Blue */
                color: white;
            }
        </style>
    """, unsafe_allow_html=True)

    st.title("Welcome to Umar-AI-Presentation")


    # Input fields
    gemini_api_key = st.text_input("Enter your Gemini API Key:", type="password")
    user_topic = st.text_input("Enter your Presentation Topic:")
    num_slides = st.selectbox("Select Number of Slides:", [2, 5, 8, 10, 12], index=1)
   
    # Ensure that the button is always visible
    st.write("Click for Download")
    if st.button("Generate Presentation"):
        if not gemini_api_key or not user_topic:
            st.error("Please provide both Gemini API Key and Topic.")
        else:
            # Fetch images related to the topic
            st.write("Fetching images...")
            image_urls = fetch_unsplash_images(user_topic, count=5)
            if not image_urls:
                st.warning("No images found for the topic.")

            # Generate presentation content
            st.write("Generating presentation using AI...")
            try:
                generated_content = generate_slide_content(user_topic, gemini_api_key, num_slides)
                st.write("Creating PowerPoint presentation...")
                ppt_filename = create_ppt(user_topic, generated_content, image_urls)
                st.success(f"Presentation '{ppt_filename}' created successfully!")

                # Provide download link for the generated PowerPoint
                with open(ppt_filename, "rb") as file:
                    st.write("Click for Download")
                    st.download_button(label="Download PowerPoint", data=file, file_name=ppt_filename, mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
            except Exception as e:
                st.error(f"An error occurred: {e}")

# Run the Streamlit app
if __name__ == "__main__":
    app()
